from langchain_core.tools import tool


@tool
def read_document(file_path: str, page_range: str = "") -> str:
    """Extract text content from a document file.

    Supports PDF, DOCX, XLSX, PPTX, HTML, and plain text files.
    For PDFs, optionally specify page_range like '1-5' or '3,7,10'.

    Args:
        file_path: Path to the document file
        page_range: Optional page range for PDFs (e.g. '1-5' or '3,7,10')

    Returns:
        Extracted text content from the document
    """
    import os

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return _read_pdf(file_path, page_range)
    elif ext == ".docx":
        return _read_docx(file_path)
    elif ext == ".xlsx":
        return _read_xlsx(file_path)
    elif ext == ".pptx":
        return _read_pptx(file_path)
    elif ext in (".html", ".htm"):
        return _read_html(file_path)
    elif ext in (".md", ".txt", ".rst", ".tex"):
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    else:
        return f"Unsupported file format: {ext}"


def _read_pdf(file_path: str, page_range: str = "") -> str:
    import fitz

    doc = fitz.open(file_path)
    pages = _parse_page_range(page_range, len(doc))
    text_parts = []
    for page_num in pages:
        page = doc[page_num]
        text_parts.append(page.get_text())
    doc.close()
    return "\n\n".join(text_parts)


def _parse_page_range(page_range: str, total_pages: int) -> list[int]:
    if not page_range:
        return list(range(total_pages))
    pages = set()
    for part in page_range.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            start = max(0, int(start) - 1)
            end = min(total_pages, int(end))
            pages.update(range(start, end))
        else:
            p = int(part) - 1
            if 0 <= p < total_pages:
                pages.add(p)
    return sorted(pages)


def _read_docx(file_path: str) -> str:
    from docx import Document

    doc = Document(file_path)
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        paragraphs.append("\n".join(rows))
    return "\n\n".join(paragraphs)


def _read_xlsx(file_path: str) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    output = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        output.append(f"## Sheet: {sheet_name}")
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append(" | ".join(cells))
        output.append("\n".join(rows))
    wb.close()
    return "\n\n".join(output)


def _read_pptx(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_parts.append(para.text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_parts.append(" | ".join(cells))
        slides_text.append("\n".join(slide_parts))
    return "\n\n".join(slides_text)


def _read_html(file_path: str) -> str:
    from bs4 import BeautifulSoup

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)
